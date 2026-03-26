"""
File Viewer Views - Preview files in browser without downloading
"""
import os
import mimetypes
from django.contrib.auth.decorators import login_required
from django.shortcuts import get_object_or_404
from django.http import HttpResponse, FileResponse, Http404, JsonResponse
from django.views.decorators.http import require_GET
from accounts.models import WorkItemAttachment


@login_required
@require_GET
def preview_file(request, attachment_id):
    """
    Serve file for preview with appropriate content type.
    Supports inline viewing for PDFs, images, and text files.
    Works with both local storage and Cloudinary.
    """
    from django.shortcuts import redirect
    from django.conf import settings
    
    attachment = get_object_or_404(WorkItemAttachment, id=attachment_id)
    
    if not attachment.file:
        raise Http404("File not found.")
    
    try:
        file_name = os.path.basename(attachment.file.name)
        file_ext = os.path.splitext(file_name)[1].lower()
        
        # Check if using Cloudinary (file has .url but no .path)
        try:
            # Try to get local path
            file_path = attachment.file.path
            is_cloudinary = False
        except (AttributeError, NotImplementedError):
            # File is on Cloudinary
            is_cloudinary = True
        
        if is_cloudinary:
            # For Cloudinary files, redirect to the Cloudinary URL
            # Cloudinary automatically serves files with correct content-type
            cloudinary_url = attachment.file.url
            
            # For documents, add fl_attachment flag to force inline viewing
            if file_ext in ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt']:
                # Add transformation to force inline display
                if '?' in cloudinary_url:
                    cloudinary_url += '&fl_attachment'
                else:
                    cloudinary_url += '?fl_attachment'
            
            return redirect(cloudinary_url)
        else:
            # Local file handling (for development)
            # Determine content type
            content_type, _ = mimetypes.guess_type(file_name)
            if not content_type:
                content_type = 'application/octet-stream'
            
            # Open file
            file_handle = open(file_path, 'rb')
            
            # For PDFs and images, use inline disposition
            if file_ext in ['.pdf', '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp']:
                response = FileResponse(file_handle, content_type=content_type)
                response['Content-Disposition'] = f'inline; filename="{file_name}"'
            else:
                # For other files, still serve inline but let browser decide
                response = FileResponse(file_handle, content_type=content_type)
                response['Content-Disposition'] = f'inline; filename="{file_name}"'
            
            return response
        
    except Exception as e:
        raise Http404(f"Error serving file: {str(e)}")


@login_required
@require_GET
def convert_docx_to_html(request, attachment_id):
    """
    Convert DOCX file to HTML for preview.
    Works with both local storage and Cloudinary.
    Requires python-docx library: pip install python-docx
    """
    import logging
    import base64
    import tempfile
    import requests
    logger = logging.getLogger(__name__)
    
    attachment = get_object_or_404(WorkItemAttachment, id=attachment_id)
    
    if not attachment.file:
        return JsonResponse({"error": "File not found"}, status=404)
    
    file_name = os.path.basename(attachment.file.name)
    file_ext = os.path.splitext(file_name)[1].lower()
    
    if file_ext not in ['.docx', '.doc']:
        return JsonResponse({"error": "Not a Word document"}, status=400)
    
    try:
        # Try to import python-docx
        try:
            from docx import Document
            from docx.oxml.text.paragraph import CT_P
            from docx.oxml.table import CT_Tbl
            from docx.table import _Cell, Table
            from docx.text.paragraph import Paragraph
        except ImportError:
            return JsonResponse({
                "success": False,
                "error": "python-docx not installed",
                "message": "To enable DOCX preview, install: pip install python-docx"
            }, status=200)
        
        # Only works with .docx (not .doc)
        if file_ext == '.doc':
            return JsonResponse({
                "success": False,
                "error": "Legacy .doc format not supported",
                "message": "Only .docx files can be previewed. Please download to view."
            }, status=200)
        
        # Check if file is on Cloudinary or local
        try:
            file_path = attachment.file.path
            is_cloudinary = False
        except (AttributeError, NotImplementedError):
            is_cloudinary = True
        
        # Load document
        try:
            if is_cloudinary:
                # Download from Cloudinary to temporary file
                cloudinary_url = attachment.file.url
                response = requests.get(cloudinary_url, timeout=30)
                response.raise_for_status()
                
                # Create temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                    tmp_file.write(response.content)
                    tmp_file_path = tmp_file.name
                
                try:
                    with open(tmp_file_path, 'rb') as docx_file:
                        doc = Document(docx_file)
                finally:
                    # Clean up temporary file
                    try:
                        os.unlink(tmp_file_path)
                    except:
                        pass
            else:
                # Local file
                with open(file_path, 'rb') as docx_file:
                    doc = Document(docx_file)
                    
        except Exception as e:
            logger.error(f"Failed to load DOCX file: {str(e)}")
            return JsonResponse({
                "success": False,
                "error": "Failed to load document",
                "message": f"This document contains unreadable content or is corrupted. Please download to view in Microsoft Word."
            }, status=200)
        
        # Convert to HTML
        html_content = '<div style="font-family: Arial, sans-serif; line-height: 1.6; color: #1e293b; max-width: 800px; margin: 0 auto; padding: 2rem;">'
        
        has_content = False
        
        # Process paragraphs and images
        for paragraph in doc.paragraphs:
            # Check for images in paragraph
            for run in paragraph.runs:
                if 'graphicData' in run._element.xml:
                    try:
                        # Extract image
                        for rel in doc.part.rels.values():
                            if "image" in rel.target_ref:
                                image_part = rel.target_part
                                image_bytes = image_part.blob
                                image_base64 = base64.b64encode(image_bytes).decode()
                                content_type = image_part.content_type
                                html_content += f'<img src="data:{content_type};base64,{image_base64}" style="max-width: 100%; height: auto; margin: 1rem 0; border-radius: 8px; box-shadow: 0 2px 8px rgba(0,0,0,0.1);" />'
                                has_content = True
                    except:
                        pass
            
            if not paragraph.text.strip():
                continue
            
            has_content = True
            
            # Check if it's a heading
            if paragraph.style.name.startswith('Heading'):
                level = paragraph.style.name.replace('Heading ', '')
                try:
                    level = int(level)
                    html_content += f'<h{level} style="margin-top: 1.5rem; margin-bottom: 0.75rem; color: #0f172a;">{paragraph.text}</h{level}>'
                except:
                    html_content += f'<p style="margin-bottom: 1rem;">{paragraph.text}</p>'
            else:
                # Regular paragraph
                text = paragraph.text
                
                # Apply basic formatting
                for run in paragraph.runs:
                    if run.bold and run.text:
                        text = text.replace(run.text, f'<strong>{run.text}</strong>')
                    if run.italic and run.text:
                        text = text.replace(run.text, f'<em>{run.text}</em>')
                    if run.underline and run.text:
                        text = text.replace(run.text, f'<u>{run.text}</u>')
                
                html_content += f'<p style="margin-bottom: 1rem;">{text}</p>'
        
        # Add tables
        for table in doc.tables:
            has_content = True
            html_content += '<table style="width: 100%; border-collapse: collapse; margin: 1.5rem 0; border: 1px solid #e2e8f0;">'
            for row in table.rows:
                html_content += '<tr>'
                for cell in row.cells:
                    html_content += f'<td style="padding: 0.75rem; border: 1px solid #e2e8f0; background: #f8fafc;">{cell.text}</td>'
                html_content += '</tr>'
            html_content += '</table>'
        
        html_content += '</div>'
        
        if not has_content:
            return JsonResponse({
                "success": False,
                "error": "No content found",
                "message": "This document appears to be empty or contains only images that couldn't be extracted. Please download to view."
            }, status=200)
        
        return JsonResponse({
            "success": True,
            "html": html_content
        })
        
    except Exception as e:
        logger.error(f"DOCX conversion error: {str(e)}", exc_info=True)
        return JsonResponse({
            "success": False,
            "error": "Conversion failed",
            "message": "This document contains complex formatting or embedded content that cannot be previewed. Please download to view in Microsoft Word."
        }, status=200)


@login_required
@require_GET
def convert_excel_to_html(request, attachment_id):
    """
    Convert Excel file to HTML table for preview.
    Works with both local storage and Cloudinary.
    Requires openpyxl library: pip install openpyxl
    """
    import tempfile
    import requests
    
    attachment = get_object_or_404(WorkItemAttachment, id=attachment_id)
    
    if not attachment.file:
        return JsonResponse({"error": "File not found"}, status=404)
    
    file_name = os.path.basename(attachment.file.name)
    file_ext = os.path.splitext(file_name)[1].lower()
    
    if file_ext not in ['.xlsx', '.xls']:
        return JsonResponse({"error": "Not an Excel document"}, status=400)
    
    try:
        # Try to import openpyxl
        try:
            from openpyxl import load_workbook
        except ImportError:
            return JsonResponse({
                "error": "openpyxl not installed",
                "message": "To enable Excel preview, install: pip install openpyxl"
            }, status=500)
        
        # Only works with .xlsx (not .xls)
        if file_ext == '.xls':
            return JsonResponse({
                "error": "Legacy .xls format not supported",
                "message": "Only .xlsx files can be previewed. Please download to view."
            }, status=400)
        
        # Check if file is on Cloudinary or local
        try:
            file_path = attachment.file.path
            is_cloudinary = False
        except (AttributeError, NotImplementedError):
            is_cloudinary = True
        
        # Load workbook
        if is_cloudinary:
            # Download from Cloudinary to temporary file
            cloudinary_url = attachment.file.url
            response = requests.get(cloudinary_url, timeout=30)
            response.raise_for_status()
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_file:
                tmp_file.write(response.content)
                tmp_file_path = tmp_file.name
            
            try:
                wb = load_workbook(tmp_file_path, data_only=True)
            finally:
                # Clean up temporary file
                try:
                    os.unlink(tmp_file_path)
                except:
                    pass
        else:
            # Local file
            wb = load_workbook(file_path, data_only=True)
        
        # Convert all sheets to HTML
        html_content = '<div style="font-family: Arial, sans-serif; padding: 2rem; max-width: 100%; overflow-x: auto;">'
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            
            # Add sheet name as heading
            if len(wb.sheetnames) > 1:
                html_content += f'<h3 style="margin-top: 2rem; margin-bottom: 1rem; color: #0f172a; border-bottom: 2px solid #e2e8f0; padding-bottom: 0.5rem;">{sheet_name}</h3>'
            
            # Start table
            html_content += '<table style="width: 100%; border-collapse: collapse; margin-bottom: 2rem; font-size: 0.875rem; box-shadow: 0 1px 3px rgba(0,0,0,0.1);">'
            
            # Get max row and column
            max_row = sheet.max_row
            max_col = sheet.max_column
            
            # Limit preview to first 100 rows
            preview_rows = min(max_row, 100)
            
            for row_idx, row in enumerate(sheet.iter_rows(max_row=preview_rows, values_only=True), 1):
                html_content += '<tr>'
                
                for col_idx, cell_value in enumerate(row, 1):
                    # Style first row as header
                    if row_idx == 1:
                        html_content += f'<th style="padding: 0.75rem; border: 1px solid #cbd5e1; background: #f1f5f9; text-align: left; font-weight: 600; color: #1e293b;">{cell_value if cell_value is not None else ""}</th>'
                    else:
                        html_content += f'<td style="padding: 0.75rem; border: 1px solid #e2e8f0; background: white;">{cell_value if cell_value is not None else ""}</td>'
                
                html_content += '</tr>'
            
            html_content += '</table>'
            
            # Add note if rows were truncated
            if max_row > 100:
                html_content += f'<p style="color: #64748b; font-style: italic; margin-bottom: 2rem;">Showing first 100 of {max_row} rows. Download file to view all data.</p>'
        
        html_content += '</div>'
        
        return JsonResponse({
            "success": True,
            "html": html_content
        })
        
    except Exception as e:
        return JsonResponse({
            "error": "Conversion failed",
            "message": str(e)
        }, status=500)


@login_required
@require_GET
def convert_pptx_to_html(request, attachment_id):
    """
    Convert PowerPoint file to HTML for preview.
    Works with both local storage and Cloudinary.
    Requires python-pptx library: pip install python-pptx
    """
    import tempfile
    import requests
    
    attachment = get_object_or_404(WorkItemAttachment, id=attachment_id)
    
    if not attachment.file:
        return JsonResponse({"error": "File not found"}, status=404)
    
    file_name = os.path.basename(attachment.file.name)
    file_ext = os.path.splitext(file_name)[1].lower()
    
    if file_ext not in ['.pptx', '.ppt']:
        return JsonResponse({"error": "Not a PowerPoint document"}, status=400)
    
    try:
        # Try to import python-pptx
        try:
            from pptx import Presentation
        except ImportError:
            return JsonResponse({
                "error": "python-pptx not installed",
                "message": "To enable PowerPoint preview, install: pip install python-pptx"
            }, status=500)
        
        # Only works with .pptx (not .ppt)
        if file_ext == '.ppt':
            return JsonResponse({
                "error": "Legacy .ppt format not supported",
                "message": "Only .pptx files can be previewed. Please download to view."
            }, status=400)
        
        # Check if file is on Cloudinary or local
        try:
            file_path = attachment.file.path
            is_cloudinary = False
        except (AttributeError, NotImplementedError):
            is_cloudinary = True
        
        # Load presentation
        if is_cloudinary:
            # Download from Cloudinary to temporary file
            cloudinary_url = attachment.file.url
            response = requests.get(cloudinary_url, timeout=30)
            response.raise_for_status()
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                tmp_file.write(response.content)
                tmp_file_path = tmp_file.name
            
            try:
                prs = Presentation(tmp_file_path)
            finally:
                # Clean up temporary file
                try:
                    os.unlink(tmp_file_path)
                except:
                    pass
        else:
            # Local file
            prs = Presentation(file_path)
        
        # Convert slides to HTML
        html_content = '<div style="font-family: Arial, sans-serif; padding: 2rem; max-width: 1200px; margin: 0 auto;">'
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            # Add slide number
            html_content += f'<div style="margin-bottom: 3rem; border: 2px solid #e2e8f0; border-radius: 12px; overflow: hidden; box-shadow: 0 4px 6px rgba(0,0,0,0.1);">'
            html_content += f'<div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 1rem 1.5rem; font-weight: 600; font-size: 1.1rem;">Slide {slide_idx}</div>'
            html_content += '<div style="background: white; padding: 2rem;">'
            
            # Extract text from shapes
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Check if it's a title (usually first text or larger)
                    if slide_idx == 1 and len(slide_text) == 0:
                        html_content += f'<h2 style="color: #1e293b; margin-bottom: 1.5rem; font-size: 1.75rem;">{text}</h2>'
                    elif len(slide_text) == 0:
                        html_content += f'<h3 style="color: #334155; margin-bottom: 1rem; font-size: 1.5rem;">{text}</h3>'
                    else:
                        # Check if it's a bullet point
                        if text.startswith('•') or text.startswith('-'):
                            html_content += f'<li style="margin-bottom: 0.5rem; color: #475569; line-height: 1.6;">{text.lstrip("•- ")}</li>'
                        else:
                            html_content += f'<p style="margin-bottom: 1rem; color: #475569; line-height: 1.6;">{text}</p>'
                    
                    slide_text.append(text)
            
            if not slide_text:
                html_content += '<p style="color: #94a3b8; font-style: italic;">No text content on this slide</p>'
            
            html_content += '</div></div>'
        
        html_content += f'<div style="margin-top: 2rem; padding: 1rem; background: #f1f5f9; border-radius: 8px; text-align: center; color: #64748b;">'
        html_content += f'<i class="fas fa-info-circle me-2"></i>Total Slides: {len(prs.slides)}'
        html_content += '</div></div>'
        
        return JsonResponse({
            "success": True,
            "html": html_content
        })
        
    except Exception as e:
        return JsonResponse({
            "error": "Conversion failed",
            "message": str(e)
        }, status=500)


@login_required
@require_GET
def get_file_info(request, attachment_id):
    """
    Get file metadata for preview modal.
    Returns JSON with file information.
    """
    attachment = get_object_or_404(WorkItemAttachment, id=attachment_id)
    
    if not attachment.file:
        return JsonResponse({"error": "File not found"}, status=404)
    
    file_name = os.path.basename(attachment.file.name)
    file_ext = os.path.splitext(file_name)[1].lower()
    file_size = attachment.file.size
    
    # Determine file type category
    file_type = 'unknown'
    if file_ext in ['.pdf']:
        file_type = 'pdf'
    elif file_ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp']:
        file_type = 'image'
    elif file_ext in ['.doc', '.docx']:
        file_type = 'word'
    elif file_ext in ['.xls', '.xlsx']:
        file_type = 'excel'
    elif file_ext in ['.ppt', '.pptx']:
        file_type = 'powerpoint'
    elif file_ext in ['.txt', '.csv', '.log', '.md']:
        file_type = 'text'
    elif file_ext in ['.zip', '.rar', '.7z', '.tar', '.gz']:
        file_type = 'archive'
    
    # Detect URL prefix based on request path
    url_prefix = '/user' if request.path.startswith('/user/') else '/admin'
    
    return JsonResponse({
        'id': attachment.id,
        'name': file_name,
        'extension': file_ext,
        'size': file_size,
        'type': file_type,
        'url': f'{url_prefix}/documents/files/preview/{attachment.id}/',
        'download_url': f'{url_prefix}/documents/download/{attachment.id}/',
        'convert_url': (
            f'{url_prefix}/documents/files/convert-docx/{attachment.id}/' if file_type == 'word' else
            f'{url_prefix}/documents/files/convert-excel/{attachment.id}/' if file_type == 'excel' else
            f'{url_prefix}/documents/files/convert-pptx/{attachment.id}/' if file_type == 'powerpoint' else
            None
        ),
    })
